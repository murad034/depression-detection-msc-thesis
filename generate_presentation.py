"""
Generate a professional MSc thesis defense PowerPoint presentation.
BEHAVIORAL AND ACADEMIC INDICATOR-BASED DEPRESSION SCREENING
AMONG HIGHER EDUCATION STUDENTS: A COMPARATIVE MACHINE LEARNING STUDY
Author: Md. Murad Hossain
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
VIS_DIR  = os.path.join(BASE_DIR, "visualizations")

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x3E, 0x6E)   # dark navy  — slide header bg
BLUE    = RGBColor(0x2E, 0x74, 0xB5)   # medium blue — accent
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)   # light gray  — body bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
GREEN   = RGBColor(0x37, 0x56, 0x23)
RED     = RGBColor(0xC0, 0x00, 0x00)
ACCENT  = RGBColor(0xED, 0x7D, 0x31)   # orange accent

TITLE_FONT  = "Calibri"
BODY_FONT   = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_name=BODY_FONT, size=18, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, wrap=True, valign=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_content_box(slide, lines, l, t, w, h,
                    size=16, color=BLACK, bullet=True, spacing=1.15):
    """Add a list of lines as bullet points in a textbox."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after  = Pt(3)
        if bullet and not line.startswith("   "):
            p.level = 0
        run = p.add_run()
        run.text = line
        run.font.name  = BODY_FONT
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb

def slide_header(slide, title_text, subtitle=None):
    """Navy top banner + white title text."""
    add_rect(slide, 0, 0, 13.33, 1.25, fill=NAVY)
    add_text(slide, title_text,
             0.35, 0.12, 12.5, 0.9,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_rect(slide, 0, 1.25, 13.33, 0.32, fill=BLUE)
        add_text(slide, subtitle,
                 0.35, 1.28, 12.5, 0.28,
                 size=14, color=WHITE, italic=True)
    add_rect(slide, 0, 1.25 if not subtitle else 1.57, 13.33, 0.04, fill=ACCENT)

def footer(slide, num, total):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=NAVY)
    add_text(slide, "Md. Murad Hossain  |  MSc Thesis Defense  |  Dept. of CSE, Gopalganj Science & Technology University",
             0.3, 7.22, 11.5, 0.26, size=9, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f"{num} / {total}",
             12.8, 7.22, 0.5, 0.26, size=9, color=WHITE, align=PP_ALIGN.RIGHT)

def add_image_safe(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

TOTAL = 14

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 2.8, 13.33, 0.06, fill=ACCENT)
add_rect(sl, 0, 4.55, 13.33, 0.06, fill=ACCENT)

add_text(sl, "MSc Thesis Presentation",
         1, 0.5, 11.33, 0.6, size=16, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
add_text(sl,
         "BEHAVIORAL AND ACADEMIC INDICATOR-BASED\nDEPRESSION SCREENING AMONG HIGHER\nEDUCATION STUDENTS",
         0.8, 1.1, 11.73, 1.8, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "A COMPARATIVE MACHINE LEARNING STUDY",
         1, 2.88, 11.33, 0.55, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, "Submitted by:", 3.5, 3.65, 2.5, 0.4, size=13, color=WHITE)
add_text(sl, "Md. Murad Hossain", 3.5, 4.0, 6, 0.5, size=20, bold=True, color=WHITE)
add_text(sl, "Supervised by:", 3.5, 4.5, 2.5, 0.4, size=13, color=WHITE)
add_text(sl, "Dr. Mrinal Kanti Baowaly", 3.5, 4.85, 6, 0.5, size=18, color=WHITE)
add_text(sl, "Department of Computer Science & Engineering\nGopalganj Science and Technology University, Bangladesh",
         2.5, 5.45, 8.33, 0.7, size=13, italic=True, color=RGBColor(0xBF, 0xD7, 0xED), align=PP_ALIGN.CENTER)
add_text(sl, "March 2026",
         5.5, 6.3, 2.33, 0.4, size=12, color=RGBColor(0xBF, 0xD7, 0xED), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE / AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Presentation Outline")
footer(sl, 2, TOTAL)

items = [
    ("01", "Problem Statement & Motivation"),
    ("02", "Research Objectives"),
    ("03", "Literature Review"),
    ("04", "Dataset & Feature Description"),
    ("05", "Methodology & Preprocessing"),
    ("06", "Classification Algorithms"),
    ("07", "Experimental Setup"),
    ("08", "Results & Performance Analysis"),
    ("09", "Conclusion & Future Work"),
]

col_w = 6.0
for i, (num, title) in enumerate(items):
    row = i % 5; col = i // 5
    lx = 0.5 + col * 6.4; ty = 1.9 + row * 0.96
    add_rect(sl, lx, ty, col_w, 0.78, fill=NAVY if i % 2 == 0 else BLUE)
    add_text(sl, num, lx + 0.15, ty + 0.15, 0.55, 0.5,
             size=20, bold=True, color=ACCENT)
    add_text(sl, title, lx + 0.75, ty + 0.18, col_w - 0.85, 0.5,
             size=15, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT & MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Problem Statement & Motivation",
             "Why does this research matter?")
footer(sl, 3, TOTAL)

# Left box
add_rect(sl, 0.4, 1.75, 5.9, 5.1, fill=WHITE, line=BLUE)
add_rect(sl, 0.4, 1.75, 5.9, 0.48, fill=NAVY)
add_text(sl, "  The Challenge", 0.4, 1.75, 5.9, 0.48, size=14, bold=True, color=WHITE)
add_content_box(sl, [
    "▸  280 million people affected by depression globally (WHO)",
    "▸  Depression prevalence: 23–41% among higher education students in South Asia",
    "▸  Clinical screening (PHQ-9, BDI) is resource-intensive and difficult to scale",
    "▸  Social stigma prevents students from seeking help",
    "▸  No systematic automated tool for early detection at institutional level",
], 0.55, 2.32, 5.65, 4.3, size=14)

# Right box
add_rect(sl, 6.9, 1.75, 5.9, 5.1, fill=WHITE, line=BLUE)
add_rect(sl, 6.9, 1.75, 5.9, 0.48, fill=NAVY)
add_text(sl, "  Bangladesh Context", 6.9, 1.75, 5.9, 0.48, size=14, bold=True, color=WHITE)
add_content_box(sl, [
    "▸  1.5+ million students enrolled in degree programs",
    "▸  Rapid lifestyle changes increasing psychological burden",
    "▸  Academic pressure, financial strain, social expectations",
    "▸  Limited mental health infrastructure at institutions",
    "▸  Urgent need: low-cost, automated, scalable screening",
], 7.05, 2.32, 5.65, 4.3, size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESEARCH OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Research Objectives")
footer(sl, 4, TOTAL)

objectives = [
    ("01", "Construct a clinically-informed behavioral & academic survey dataset from\n       539 higher education students across 23 feature dimensions"),
    ("02", "Design and implement a systematic data preprocessing pipeline\n       (encoding, imputation, and Z-score normalization)"),
    ("03", "Train and evaluate 4 supervised ML classifiers:\n       SVM, Logistic Regression, XGBoost, and MLP"),
    ("04", "Compare models using Accuracy, Precision, Recall,\n       F1-Score, and ROC AUC under 10-fold cross-validation"),
    ("05", "Identify the optimal model for real-world\n       student depression screening deployment"),
]

for i, (num, obj) in enumerate(objectives):
    ty = 1.55 + i * 1.04
    add_rect(sl, 0.5, ty, 0.7, 0.82, fill=NAVY)
    add_text(sl, num, 0.5, ty+0.12, 0.7, 0.6, size=17, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.28, ty, 11.55, 0.82, fill=WHITE, line=BLUE)
    add_text(sl, obj, 1.45, ty + 0.1, 11.2, 0.72, size=14, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LITERATURE REVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Literature Review", "Key prior work in ML-based depression detection")
footer(sl, 5, TOTAL)

papers = [
    ("Mohd & Yahya\n(2018)",     "LR + ANN on student questionnaire data\n→ 72–78% accuracy"),
    ("Deshpande & Rao\n(2017)",  "SVM + Naive Bayes on emotion-tagged social content\n→ F1-score based evaluation"),
    ("Cacheda et al.\n(2019)",   "Dual Random Forest on social network streams\n→ Early detection, dual model outperforms single"),
    ("Bae & Lee\n(2015)",        "Data mining of Korean adolescent clinical data\n→ Depression = strongest predictor of suicide"),
    ("Dipnall et al.\n(2016)",   "Fused ML + statistics for biomarker discovery\n→ Ensemble learning effective in psychiatry"),
    ("Shrestha\n(2018)",         "Twitter text features for depression prediction\n→ Limited generalizability due to small dataset"),
]

for i, (study, finding) in enumerate(papers):
    row = i % 3; col = i // 3
    lx = 0.4 + col * 6.45; ty = 1.75 + row * 1.7
    add_rect(sl, lx, ty, 6.1, 1.5, fill=WHITE, line=BLUE)
    add_rect(sl, lx, ty, 2.1, 1.5, fill=NAVY)
    add_text(sl, study, lx + 0.1, ty + 0.28, 1.92, 0.95,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, finding, lx + 2.22, ty + 0.15, 3.7, 1.22, size=12, color=BLACK)

add_rect(sl, 0.4, 6.95, 12.53, 0.15, fill=ACCENT)
add_text(sl, "Research Gap: No multi-classifier benchmarking framework on behavioral survey data for Bangladeshi higher education students",
         0.5, 6.93, 12.3, 0.3, size=11, italic=True, color=NAVY, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATASET & FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Dataset & Feature Description",
             "539 students · 23 features · binary class label")
footer(sl, 6, TOTAL)

# Stats boxes
stats = [("539", "Student Records"), ("23", "Feature Dimensions"),
         ("40.62%", "Depressed"), ("59.38%", "Not Depressed")]
for i, (val, label) in enumerate(stats):
    lx = 0.35 + i * 3.15
    add_rect(sl, lx, 1.72, 2.8, 1.05, fill=NAVY)
    add_text(sl, val,   lx+0.1, 1.73, 2.6, 0.62, size=26, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, label, lx+0.1, 2.32, 2.6, 0.42, size=13, color=WHITE,
             align=PP_ALIGN.CENTER)

# Feature categories
cats = [
    ("Behavioral /\nPsychological",
     ["Social norms acceptance", "Personality type", "Solitude comfort",
      "Talking about problems", "Hanging out with friends"]),
    ("Academic /\nPerformance",
     ["SSC Result", "HSC Result", "University CGPA",
      "Challenging Education System", "Extracurricular activities"]),
    ("Digital /\nLifestyle",
     ["Smartphone ownership", "Daily sleep duration",
      "Contentment in current role", "Work anxiety"]),
    ("Family /\nSocial",
     ["Family contentment", "Feel like a burden",
      "Comfortable environment", "Family size"]),
    ("Clinical /\nRisk",
     ["Suicide attempt history", "Thoughts on suicide",
      "Age (17–33)", "Gender", "  → Class Label: Depression (Yes/No)"]),
]

bx = 0.35
for cat, items in cats:
    add_rect(sl, bx, 2.98, 2.42, 0.42, fill=BLUE)
    add_text(sl, cat, bx+0.05, 2.99, 2.32, 0.4, size=11, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_rect(sl, bx, 3.42 + j * 0.62, 2.42, 0.58,
                 fill=WHITE if j % 2 == 0 else RGBColor(0xE8, 0xF0, 0xF8))
        add_text(sl, item, bx+0.1, 3.45 + j * 0.62, 2.28, 0.52, size=10.5, color=BLACK)
    bx += 2.52

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — METHODOLOGY & SYSTEM WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Methodology & System Workflow")
footer(sl, 7, TOTAL)

add_image_safe(sl, os.path.join(VIS_DIR, "fig_page1_img.png"),
               0.4, 1.55, 12.5, 2.1)
add_text(sl, "Fig. 1: End-to-End System Pipeline — from Survey Design to Depression Risk Output",
         1, 3.68, 11.33, 0.35, size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# Preprocessing steps
add_rect(sl, 0.35, 4.12, 12.63, 0.38, fill=NAVY)
add_text(sl, "  Preprocessing Pipeline", 0.35, 4.12, 12.63, 0.38,
         size=14, bold=True, color=WHITE)

steps = [
    ("1. Binary\nEncoding", "Yes/No → 1/0\n(13 features)"),
    ("2. One-Hot\nEncoding", "Gender, Personality\nEnvironment → binary"),
    ("3. Median\nImputation", "Missing values\nreplaced by median"),
    ("4. Z-Score\nScaling", "μ=0, σ=1\nStandardScaler"),
    ("5. Feature\nSpace", "23 attrs →\n26 dimensions"),
]

for i, (step, detail) in enumerate(steps):
    lx = 0.4 + i * 2.51
    add_rect(sl, lx, 4.58, 2.3, 2.52, fill=WHITE, line=BLUE)
    add_rect(sl, lx, 4.58, 2.3, 0.58, fill=BLUE)
    add_text(sl, step,   lx+0.08, 4.6, 2.15, 0.55, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, detail, lx+0.08, 5.2, 2.15, 1.8,  size=12, color=BLACK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLASSIFICATION ALGORITHMS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Classification Algorithms",
             "Four complementary ML paradigms for comparative analysis")
footer(sl, 8, TOTAL)

algos = [
    ("SVM",
     RGBColor(0x1F, 0x3E, 0x6E),
     ["Kernel-based margin maximization",
      "RBF kernel for non-linear boundaries",
      "C=10, gamma='scale'",
      "Strong on small-to-medium datasets"]),
    ("Logistic\nRegression",
     BLUE,
     ["Linear probabilistic classifier",
      "Sigmoid function → probability",
      "L-BFGS solver, C=1.0, max_iter=1000",
      "Interpretable decision boundary"]),
    ("XGBoost",
     GREEN,
     ["Gradient-boosted decision trees",
      "L1/L2 regularization built-in",
      "200 estimators, max_depth=6, lr=0.1",
      "Handles missing values natively"]),
    ("MLP\n(Neural Net)",
     RED,
     ["3 hidden layers: 128 → 64 → 32",
      "ReLU activation, Adam optimizer",
      "Softmax output, max_iter=500",
      "Non-linear feature representation"]),
]

for i, (name, color, points) in enumerate(algos):
    lx = 0.4 + i * 3.22
    add_rect(sl, lx, 1.72, 2.9, 1.0, fill=color)
    add_text(sl, name, lx+0.1, 1.78, 2.7, 0.88,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, 2.72, 2.9, 4.18, fill=WHITE, line=color)
    for j, pt in enumerate(points):
        add_text(sl, f"▸  {pt}", lx+0.15, 2.85 + j*0.95, 2.65, 0.85,
                 size=12.5, color=BLACK)

# Evaluation banner
add_rect(sl, 0.35, 6.95, 12.63, 0.38, fill=ACCENT)
add_text(sl, "  Evaluation: 10-Fold Stratified Cross-Validation   |   Metrics: Accuracy · Precision · Recall · F1-Score · ROC AUC",
         0.5, 6.96, 12.3, 0.35, size=13, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EXPERIMENTAL SETUP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Experimental Setup & Evaluation Protocol")
footer(sl, 9, TOTAL)

# Left: environment
add_rect(sl, 0.4, 1.68, 5.9, 5.35, fill=WHITE, line=BLUE)
add_rect(sl, 0.4, 1.68, 5.9, 0.5, fill=NAVY)
add_text(sl, "  Simulation Environment", 0.4, 1.68, 5.9, 0.5, size=14, bold=True, color=WHITE)
env_items = [
    "Programming Language:  Python 3.9",
    "scikit-learn 1.2.2",
    "XGBoost 1.7.5",
    "pandas 1.5.3  |  NumPy 1.24.3",
    "Matplotlib 3.7.1",
]
for j, item in enumerate(env_items):
    bg = RGBColor(0xE8, 0xF0, 0xF8) if j % 2 == 0 else WHITE
    add_rect(sl, 0.4, 2.26 + j * 0.88, 5.9, 0.85, fill=bg)
    add_text(sl, f"  ✓  {item}", 0.5, 2.3 + j * 0.88, 5.7, 0.8, size=13.5, color=BLACK)

# Right: evaluation
add_rect(sl, 7.0, 1.68, 5.9, 5.35, fill=WHITE, line=BLUE)
add_rect(sl, 7.0, 1.68, 5.9, 0.5, fill=NAVY)
add_text(sl, "  Evaluation Protocol", 7.0, 1.68, 5.9, 0.5, size=14, bold=True, color=WHITE)
eval_items = [
    ("10-Fold Stratified CV", "Class ratio preserved in every fold"),
    ("Train / Test Split",    "80% (431) training  /  20% (108) testing"),
    ("Accuracy",              "(TP+TN) / (TP+TN+FP+FN)"),
    ("Precision",             "TP / (TP + FP)"),
    ("Recall (TPR)",          "TP / (TP + FN)"),
    ("F1-Score",              "2 × (Precision × Recall) / (P+R)"),
    ("ROC AUC",               "Area under ROC curve (0–1)"),
]
for j, (metric, formula) in enumerate(eval_items):
    bg = RGBColor(0xE8, 0xF0, 0xF8) if j % 2 == 0 else WHITE
    add_rect(sl, 7.0, 2.26 + j * 0.68, 5.9, 0.65, fill=bg)
    add_text(sl, metric,  7.08, 2.29 + j * 0.68, 2.3, 0.58,
             size=12, color=NAVY, bold=True)
    add_text(sl, formula, 9.42, 2.29 + j * 0.68, 3.4, 0.58,
             size=12, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS TABLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Experimental Results",
             "TABLE: Performance Comparison — 10-Fold Stratified Cross-Validation")
footer(sl, 10, TOTAL)

headers = ["Metric", "MLP", "Logistic Regression", "XGBoost", "SVM ★"]
col_ws  = [3.5, 2.0, 2.7, 2.0, 2.05]
lx_start = 0.55

# header row
lx = lx_start
for ci, (h, cw) in enumerate(zip(headers, col_ws)):
    add_rect(sl, lx, 1.66, cw, 0.55, fill=NAVY)
    add_text(sl, h, lx+0.06, 1.68, cw-0.1, 0.52,
             size=14, bold=True, color=ACCENT if ci == 4 else WHITE,
             align=PP_ALIGN.CENTER)
    lx += cw + 0.03

data = [
    ("Accuracy (%)",       "96.84", "97.77",  "97.59",  "97.96"),
    ("Avg. Precision (%)", "95.47", "97.23",  "96.30",  "97.23"),
    ("Avg. Recall (%)",    "95.37", "97.22",  "96.30",  "97.22"),
    ("Avg. F1-Score (%)",  "95.34", "97.22",  "96.30",  "97.22"),
    ("Avg. ROC AUC (%)",   "99.01", "99.50",  "99.57",  "99.11"),
]

for ri, row in enumerate(data):
    ty = 2.24 + ri * 0.82
    bg = RGBColor(0xEB, 0xF3, 0xFB) if ri % 2 == 0 else WHITE
    lx = lx_start
    for ci, (val, cw) in enumerate(zip(row, col_ws)):
        is_best = (ci == 4) and (ri < 4)  # SVM best for acc/prec/rec/f1
        is_best_auc = (ci == 3) and (ri == 4)  # XGBoost highest AUC
        cell_bg = RGBColor(0xC6, 0xEF, 0xCE) if (is_best or is_best_auc) else bg
        add_rect(sl, lx, ty, cw, 0.74, fill=cell_bg,
                 line=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(sl, val, lx+0.06, ty+0.1, cw-0.1, 0.55,
                 size=15, bold=(is_best or is_best_auc),
                 color=NAVY if (is_best or is_best_auc) else BLACK,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        lx += cw + 0.03

add_text(sl, "★ SVM achieved highest overall accuracy; XGBoost achieved highest ROC AUC  |  All models exceeded 96% accuracy",
         0.5, 7.02, 12.3, 0.35, size=11, italic=True, color=NAVY, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS CHARTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Results — Visual Analysis")
footer(sl, 11, TOTAL)

add_image_safe(sl, os.path.join(VIS_DIR, "model_comparison.png"),   0.3,  1.52, 6.1, 2.65)
add_image_safe(sl, os.path.join(VIS_DIR, "roc_curves.png"),         6.85, 1.52, 6.1, 2.65)
add_image_safe(sl, os.path.join(VIS_DIR, "metrics_comparison.png"), 0.3,  4.28, 6.1, 2.9)
add_image_safe(sl, os.path.join(VIS_DIR, "confusion_matrices.png"), 6.85, 4.28, 6.1, 2.9)

add_text(sl, "Fig. 3: Accuracy Comparison", 0.3, 4.1, 6.1, 0.28,
         size=10, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "Fig. 5: ROC Curves (all AUC > 99%)", 6.85, 4.1, 6.1, 0.28,
         size=10, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "Fig. 6: All Metrics Comparison", 0.3, 7.1, 6.1, 0.28,
         size=10, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "Fig. 4: Confusion Matrices", 6.85, 7.1, 6.1, 0.28,
         size=10, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY FINDINGS / DISCUSSION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Key Findings & Discussion")
footer(sl, 12, TOTAL)

findings = [
    (NAVY,  "SVM — Best Accuracy",
     "97.96% accuracy confirms that RBF kernel maximum-margin classification\ncaptures non-linear patterns in behavioral survey features most effectively."),
    (BLUE,  "Logistic Regression — Strong Linear Performance",
     "97.77% accuracy despite linear decision boundary — indicates the\npreprocessed feature space is substantially linearly separable."),
    (GREEN, "XGBoost — Highest ROC AUC",
     "99.57% ROC AUC demonstrates superior ranking capability;\ngradient boosting effectively captures feature interaction effects."),
    (RGBColor(0x7B, 0x7B, 0x7B), "MLP — Solid Neural Network Baseline",
     "96.84% accuracy; hierarchical neural feature representations\nvalid for complex non-linear depression indicator patterns."),
    (ACCENT, "Overall Conclusion",
     "All 4 models exceed 96% accuracy and 99% AUC — validating that\nbehavioral survey data + principled preprocessing = reliable depression screening."),
]

for i, (color, title, detail) in enumerate(findings):
    ty = 1.55 + i * 1.04
    add_rect(sl, 0.4, ty, 0.12, 0.9, fill=color)
    add_rect(sl, 0.58, ty, 12.35, 0.9, fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC))
    add_text(sl, title,  0.75, ty+0.04, 5.0, 0.42, size=14, bold=True, color=color)
    add_text(sl, detail, 0.75, ty+0.44, 12.0, 0.45, size=12, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
slide_header(sl, "Conclusion & Future Work")
footer(sl, 13, TOTAL)

add_rect(sl, 0.4, 1.68, 6.1, 5.12, fill=WHITE, line=BLUE)
add_rect(sl, 0.4, 1.68, 6.1, 0.5, fill=NAVY)
add_text(sl, "  Conclusion", 0.4, 1.68, 6.1, 0.5, size=15, bold=True, color=WHITE)
conc = [
    "✔  Proposed ML-based depression screening framework for\n    higher education students using 23-feature survey data",
    "✔  SVM achieved highest accuracy: 97.96%",
    "✔  All 4 models exceeded 97% accuracy and 99% ROC AUC",
    "✔  Survey-driven, non-invasive, low-cost approach —\n    deployable at institutional scale",
    "✔  SVM recommended as optimal classifier for this domain",
]
for j, c in enumerate(conc):
    bg = RGBColor(0xEB, 0xF3, 0xFB) if j % 2 == 0 else WHITE
    add_rect(sl, 0.4, 2.26 + j * 0.88, 6.1, 0.84, fill=bg)
    add_text(sl, c, 0.55, 2.3 + j * 0.88, 5.85, 0.8, size=12.5, color=BLACK)

add_rect(sl, 6.9, 1.68, 6.0, 5.12, fill=WHITE, line=BLUE)
add_rect(sl, 6.9, 1.68, 6.0, 0.5, fill=NAVY)
add_text(sl, "  Future Work", 6.9, 1.68, 6.0, 0.5, size=15, bold=True, color=WHITE)
future = [
    "➤  Expand dataset — multi-institution, longitudinal,\n    diverse socioeconomic backgrounds",
    "➤  Multi-class severity: mild / moderate / severe\n    aligned with PHQ-9 clinical scale",
    "➤  Deep learning: transformers, graph neural networks\n    for relational/temporal feature modeling",
    "➤  Passive sensing: smartphone usage, wearable\n    sleep data as complementary features",
    "➤  Real-time web application for university\n    counseling service integration",
]
for j, fw in enumerate(future):
    bg = RGBColor(0xEB, 0xF3, 0xFB) if j % 2 == 0 else WHITE
    add_rect(sl, 6.9, 2.26 + j * 0.88, 6.0, 0.84, fill=bg)
    add_text(sl, fw, 7.05, 2.3 + j * 0.88, 5.75, 0.8, size=12.5, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 3.15, 13.33, 0.08, fill=ACCENT)
add_rect(sl, 0, 4.85, 13.33, 0.08, fill=ACCENT)

add_text(sl, "Thank You", 1, 0.8, 11.33, 1.5,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Questions & Discussion",
         2, 2.35, 9.33, 0.8, size=24, color=ACCENT,
         align=PP_ALIGN.CENTER, italic=True)

add_text(sl, "Md. Murad Hossain",
         3.5, 3.38, 6.33, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "muradmd312@gmail.com",
         3.8, 3.95, 5.73, 0.5, size=14, color=RGBColor(0xBF, 0xD7, 0xED),
         align=PP_ALIGN.CENTER)
add_text(sl, "Supervised by: Dr. Mrinal Kanti Baowaly",
         3, 4.52, 7.33, 0.45, size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
         "Dept. of CSE  |  Gopalganj Science and Technology University, Bangladesh",
         2, 5.0, 9.33, 0.45, size=13, italic=True,
         color=RGBColor(0xBF, 0xD7, 0xED), align=PP_ALIGN.CENTER)

# Summary strip
add_rect(sl, 1.5, 5.65, 10.33, 1.28, fill=RGBColor(0x16, 0x2C, 0x52))
add_text(sl, "SVM: 97.96%     LR: 97.77%     XGBoost: 97.59%     MLP: 96.84%",
         1.8, 5.72, 9.73, 0.55, size=16, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(sl, "All ROC AUC > 99%  |  10-Fold Stratified Cross-Validation  |  539 Students  |  23 Features",
         1.8, 6.22, 9.73, 0.5, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(BASE_DIR, "Depression_Screening_MSc_Thesis_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
